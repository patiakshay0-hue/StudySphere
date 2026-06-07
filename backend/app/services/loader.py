"""Document ingestion: extract text from PDF / DOCX / TXT and split it into
overlapping, page-aware chunks suitable for retrieval.
"""
import io
import os
import tempfile
import zipfile
from typing import List, Optional, Tuple

import docx2txt
import pdfplumber
from fastapi import UploadFile
from pptx import Presentation

SUPPORTED = {"pdf", "docx", "pptx", "txt", "md"}

CHUNK_SIZE = 900
CHUNK_OVERLAP = 150


def _detect_kind(content: bytes, ext: str) -> Optional[str]:
    """Determine the real file type from its content (magic bytes), falling back
    to the extension. Returns 'pdf' | 'docx' | 'pptx' | 'txt' | 'ole' | None.

    This makes uploads robust to misleading names like ``notes.docx.pdf``.
    """
    head = content[:8]
    if head[:4] == b"%PDF":
        return "pdf"
    if head[:4] == b"PK\x03\x04":  # ZIP container = Office Open XML (docx/pptx/xlsx)
        try:
            names = zipfile.ZipFile(io.BytesIO(content)).namelist()
            if any(n.startswith("word/") for n in names):
                return "docx"
            if any(n.startswith("ppt/") for n in names):
                return "pptx"
            if any(n.startswith("xl/") for n in names):
                return None  # xlsx not supported
        except zipfile.BadZipFile:
            return None
        return None
    if head[:4] == b"\xD0\xCF\x11\xE0":  # legacy OLE (.doc/.ppt/.xls)
        return "ole"
    # Plain text?
    try:
        content[:4096].decode("utf-8")
        return "txt"
    except UnicodeDecodeError:
        return "txt" if ext in {"txt", "md"} else None


async def extract_pages(file: UploadFile) -> Tuple[List[Tuple[Optional[int], str]], str]:
    """Return ``([(page_no, text), ...], kind)`` for an uploaded file.

    ``page_no`` is the 1-based page/slide for PDF/PPTX; ``None`` for DOCX/TXT.
    The file type is detected from the *content*, not the filename extension.
    """
    content = await file.read()
    name = file.filename or "document"
    ext = name.lower().rsplit(".", 1)[-1] if "." in name else ""

    kind = _detect_kind(content, ext)
    if kind == "ole":
        raise ValueError(
            "Old .doc/.ppt format isn't supported. Please re-save as .docx, .pptx, or PDF and upload again."
        )
    if kind is None:
        raise ValueError(
            "Couldn't read this file. Supported types: PDF, DOCX, PPTX, TXT, MD. "
            "(Tip: if it's a Word/PowerPoint file, make sure it isn't corrupted or password-protected.)"
        )

    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{kind}") as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        if kind == "pdf":
            pages = _extract_pdf(tmp_path)
        elif kind == "pptx":
            pages = _extract_pptx(tmp_path)
        elif kind == "docx":
            pages = [(None, _extract_docx(tmp_path))]
        else:  # txt / md
            pages = [(None, content.decode("utf-8", errors="ignore"))]
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.remove(tmp_path)

    pages = [(p, t) for p, t in pages if t and t.strip()]
    if not pages:
        raise ValueError(
            "No readable text could be extracted — the file may be scanned images, empty, or password-protected."
        )
    return pages, kind


def _extract_pdf(path: str) -> List[Tuple[Optional[int], str]]:
    out: List[Tuple[Optional[int], str]] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                out.append((i, text))
    return out


def _extract_docx(path: str) -> str:
    return docx2txt.process(path) or ""


def _extract_pptx(path: str) -> List[Tuple[Optional[int], str]]:
    """Extract text per slide, treating each slide as a page."""
    out: List[Tuple[Optional[int], str]] = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        text = "\n".join(parts)
        if text.strip():
            out.append((i, text))
    return out


def chunk_pages(
    pages: List[Tuple[Optional[int], str]],
    source_name: str,
    size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> List[dict]:
    """Split extracted pages into overlapping chunks, preserving page numbers."""
    chunks: List[dict] = []
    idx = 0
    for page_no, text in pages:
        normalized = " ".join(text.split())
        start = 0
        length = len(normalized)
        while start < length:
            end = min(start + size, length)
            # Prefer to break on a sentence/word boundary near the chunk end.
            if end < length:
                window = normalized[start:end]
                cut = max(window.rfind(". "), window.rfind("? "), window.rfind("! "))
                if cut > size * 0.5:
                    end = start + cut + 1
            piece = normalized[start:end].strip()
            if piece:
                chunks.append(
                    {
                        "idx": idx,
                        "text": piece,
                        "page": page_no,
                        "source": source_name,
                    }
                )
                idx += 1
            if end >= length:
                break
            start = max(end - overlap, start + 1)
    return chunks


def full_text(pages: List[Tuple[Optional[int], str]]) -> str:
    return "\n\n".join(t for _, t in pages)
